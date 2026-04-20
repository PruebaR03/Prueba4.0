import re
import pandas as pd
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def aplicar_operacion(excel_actual: str, archivo_nombre: str, operacion_nombre: str, operaciones_config: list, excel_anterior: str = None) -> str:
    """
    Aplica una operación a un archivo Excel y retorna el conteo formateado.
    Soporta operaciones con múltiples condiciones (AND).
    Retorna formato: "= número con respecto a la semana pasada (anterior)" o simple "número" si no hay anterior.
    """
    op = None
    for o in operaciones_config:
        if o.get('nombre') == operacion_nombre:
            op = o
            break
    
    if not op:
        return f"⚠️ Operación '{operacion_nombre}' no configurada"
    
    try:
        xl = pd.ExcelFile(excel_actual)
        hoja_target = None
        for hoja in xl.sheet_names:
            if hoja.lower().replace(" ", "_") == archivo_nombre.lower().replace(" ", "_"):
                hoja_target = hoja
                break
        
        if not hoja_target:
            return f"⚠️ Hoja '{archivo_nombre}' no encontrada"
        
        df = xl.parse(hoja_target)
        df.columns = [str(c).strip().strip('"').strip("'").strip().lower() for c in df.columns]
        
        # Obtener condiciones
        condiciones = op.get('condiciones', [])
        
        conteo_actual = 0
        if not condiciones:
            # Fallback a formato antiguo (una sola condición)
            columna = op.get('columna', '').strip().lower()
            valor = op.get('valor', '').strip()
            
            if not columna or columna not in df.columns:
                return f"⚠️ Columna '{columna}' no existe"
            
            # Aplicar filtro simple
            mask = _aplicar_filtro_simple(df[columna], valor)
            conteo_actual = int(mask.sum())
        else:
            # Aplicar múltiples condiciones con AND
            mask_resultado = pd.Series([True] * len(df), index=df.index)
            
            for cond in condiciones:
                columna = cond.get('columna', '').strip().lower()
                valor = cond.get('valor', '').strip()
                
                if not columna or columna not in df.columns:
                    print(f"⚠️ Columna '{columna}' no existe en operación '{operacion_nombre}'")
                    return "0"
                
                # Aplicar filtro y hacer AND con resultado anterior
                mask_cond = _aplicar_filtro_simple(df[columna], valor)
                mask_resultado = mask_resultado & mask_cond
            
            conteo_actual = int(mask_resultado.sum())
        
        # Intentar obtener conteo anterior si se proporciona Excel anterior
        conteo_anterior = None
        if excel_anterior:
            try:
                xl_ant = pd.ExcelFile(excel_anterior)
                hoja_ant = None
                for hoja in xl_ant.sheet_names:
                    if hoja.lower().replace(" ", "_") == archivo_nombre.lower().replace(" ", "_"):
                        hoja_ant = hoja
                        break
                
                if hoja_ant:
                    df_ant = xl_ant.parse(hoja_ant)
                    df_ant.columns = [str(c).strip().strip('"').strip("'").strip().lower() for c in df_ant.columns]
                    
                    if not condiciones:
                        columna = op.get('columna', '').strip().lower()
                        valor = op.get('valor', '').strip()
                        if columna and columna in df_ant.columns:
                            mask_ant = _aplicar_filtro_simple(df_ant[columna], valor)
                            conteo_anterior = int(mask_ant.sum())
                    else:
                        mask_resultado_ant = pd.Series([True] * len(df_ant), index=df_ant.index)
                        for cond in condiciones:
                            columna = cond.get('columna', '').strip().lower()
                            valor = cond.get('valor', '').strip()
                            if columna and columna in df_ant.columns:
                                mask_cond = _aplicar_filtro_simple(df_ant[columna], valor)
                                mask_resultado_ant = mask_resultado_ant & mask_cond
                        conteo_anterior = int(mask_resultado_ant.sum())
            except Exception:
                conteo_anterior = None
        
        # Formatear resultado como comparativa
        if conteo_anterior is not None:
            # Usar formato comparativo
            from .ppt_utils import formato_cambio
            texto, color, _ = formato_cambio(conteo_actual, conteo_anterior)
            return texto
        else:
            # Sin anterior, retornar solo el número
            return str(conteo_actual)
    
    except Exception as e:
        print(f"⚠️ Error aplicando operación '{operacion_nombre}': {e}")
        return "0"

def _aplicar_filtro_simple(serie: pd.Series, valor: str) -> pd.Series:
    """
    Aplica un filtro simple a una serie según el formato del valor.
    Soporta:
    - +valor+ = Contiene
    - *valor* = NO contiene (diferente de)
    - valor = Igual exacto
    """
    serie =serie.fillna('vacio')
    if valor.startswith("+") and valor.endswith("+"):
        # +valor+ = Contiene
        valor_filtro = valor.strip("+")
        return serie.astype(str).str.contains(valor_filtro, case=False, na=False)
    elif valor.startswith("*") and valor.endswith("*"):
        # *valor* = NO contiene (diferente de)
        valor_filtro = valor.strip("*")
        # Invertir la máscara de contiene
        return ~serie.astype(str).str.contains(valor_filtro, case=False, na=False)
    else:
        # Valor exacto
        return serie.astype(str).str.strip() == valor

def reemplazar_con_formato(text_frame, texto_original: str, valores: dict, excel_actual: str = None, excel_anterior: str = None, instrucciones: list = None):
    """
    Reemplaza banderas manteniendo colores y formato original.
    Soporta banderas con operaciones: <<nombre(operacion)>>
    """
    valores_norm = {}
    for k, v in valores.items():
        k_norm = k.lower()
        k_norm = re.sub(r'[\s\-\.]+', '_', k_norm)
        valores_norm[k_norm] = v

    # Encontrar banderas y sus valores (incluyendo banderas con operaciones)
    banderas = []
    for match in re.finditer(r"<<([^>]+)>>", texto_original):
        bandera_texto = match.group(1)
        op_match = re.match(r'^([^\(]+)\(([^\)]+)\)$', bandera_texto)
        
        if op_match:
            # Bandera con operación
            nombre_archivo = op_match.group(1).strip().lower()
            nombre_operacion = op_match.group(2).strip()
            
            valor_operacion = None
            if instrucciones and excel_actual:
                for inst in instrucciones:
                    if inst.get("archivo", "").lower().replace(" ", "_") == nombre_archivo.replace(" ", "_"):
                        operaciones = inst.get("operaciones", [])
                        valor_operacion = aplicar_operacion(
                            excel_actual, 
                            inst.get("archivo", ""), 
                            nombre_operacion, 
                            operaciones,
                            excel_anterior=excel_anterior
                        )
                        break
            
            if valor_operacion:
                banderas.append({
                    'start': match.start(),
                    'end': match.end(),
                    'bandera': bandera_texto,
                    'bandera_norm': nombre_archivo.lower().replace(" ", "_"),
                    'valor': valor_operacion,
                    'color': None,
                    'es_operacion': True
                })
        else:
            # Bandera normal
            bandera = bandera_texto
            bandera_norm = re.sub(r'[\s\-\.]+', '_', bandera.lower())
            if bandera_norm in valores_norm:
                banderas.append({
                    'start': match.start(),
                    'end': match.end(),
                    'bandera': bandera,
                    'bandera_norm': bandera_norm,
                    'valor': str(valores_norm[bandera_norm]),
                    'color': valores_norm.get(f'{bandera_norm}_color'),
                    'es_operacion': False
                })

    if not banderas:
        return

    print(f"                     📝 Reemplazando {len(banderas)} bandera(s)...")

    # Limpiar y reconstruir
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER

    pos_actual = 0
    for band_info in banderas:
        # Texto antes
        if pos_actual < band_info['start']:
            run_antes = p.add_run()
            run_antes.text = texto_original[pos_actual:band_info['start']]
            run_antes.font.size = Pt(10)

        valor = band_info['valor']
        
        if band_info['es_operacion']:
            # Operación: aplicar mismo formato que banderas normales
            m = re.match(r'^\s*([=↑↓]\s*\d+)\s*(.*)$', valor)
            
            if m:
                parte_color = m.group(1)
                parte_texto = m.group(2)

                simbolo = parte_color[0]
                numero = parte_color[1:]
            else:
                parte_color = valor
                parte_texto = ""

                simbolo = parte_color[0]
                numero = parte_color[1:]
            
            # Determinar color basado en símbolo Y contexto del texto
            color_op = _determinar_color_operacion(parte_color, parte_texto)
            
            # Run 1: símbolo + número (CON COLOR)
            #run_num = p.add_run()
            #run_num.text = parte_color
            #run_num.font.bold = True
            #run_num.font.size = Pt(10)
            #if color_op:
            #    run_num.font.color.rgb = color_op

            run_simbolo = p.add_run()
            run_simbolo.text = simbolo
            run_simbolo.font.bold = True
            run_simbolo.font.size = Pt(10)
            color_simbolo = _determinar_color_simbolo(simbolo)

            if color_simbolo:
                run_simbolo.font.color.rgb = color_simbolo

            run_numero = p.add_run()
            run_numero.text = numero
            run_numero.font.bold = True
            run_numero.font.size = Pt(10)
            color_numero = _determinar_color_operacion(parte_color, parte_texto)
            
            if color_numero:
                run_numero.font.color.rgb = color_numero

            if parte_texto:
                run_desc = p.add_run()
                run_desc.text = " " + parte_texto
                run_desc.font.bold = False
                run_desc.font.size = Pt(9)
        else:
            # Bandera normal: separar símbolo+número de descripción
            m = re.match(r'^\s*([=↑↓]\s*\d+)\s*(.*)$', valor)
            
            if m:
                parte_color = m.group(1)
                parte_texto = m.group(2)

                simbolo = parte_color[0]
                numero = parte_color[1:]
            else:
                parte_color = valor
                parte_texto = ""

                simbolo = parte_color[0]
                numero = parte_color[1:]
            
            run_simbolo = p.add_run()
            run_simbolo.text = simbolo
            run_simbolo.font.bold = True
            run_simbolo.font.size = Pt(10)
            color_simbolo = _determinar_color_simbolo(simbolo)

            if color_simbolo:
                run_simbolo.font.color.rgb = color_simbolo

            run_numero = p.add_run()
            run_numero.text = numero
            run_numero.font.bold = True
            run_numero.font.size = Pt(10)
            color_numero = _determinar_color_operacion(parte_color, parte_texto)
            
            if color_numero:
                run_numero.font.color.rgb = color_numero

            if parte_texto:
                run_desc = p.add_run()
                run_desc.text = " " + parte_texto
                run_desc.font.bold = False
                run_desc.font.size = Pt(9)

        print(f"                        ✓ <<{band_info['bandera']}>> → {valor}")
        
        pos_actual = band_info['end']

    # Texto restante
    if pos_actual < len(texto_original):
        run_restante = p.add_run()
        run_restante.text = texto_original[pos_actual:]
        run_restante.font.size = Pt(10)

def _determinar_color_simbolo(simbolo: str):
    if simbolo == "↑":
        return RGBColor(192, 0, 0)
    if simbolo == "↓":
        return RGBColor(0, 128, 0)
    if simbolo == "=":
        return RGBColor(0, 112, 192)
    
    return None

def _determinar_color_operacion(parte_color: str, parte_texto: str = "") -> RGBColor:
    """
    Determina el color basado en el símbolo y el contexto.
    Analiza el texto descriptivo para diferenciar entre aumento, disminución y sin cambio.
    """
    texto_completo = (parte_color + " " + parte_texto).lower()
    
    #Colores:
    #Verde: 0, 128, 0
    #Rojo:  192, 0, 0
    #Azul:  0, 112, 192
    try:
        valor = int(parte_color[1:])
    except:
        return None
    
    if valor == 0:
        return RGBColor(0, 112, 192)

    # Detectar por símbolo principal
    if "↓" in parte_color:
        # Flecha hacia abajo = Disminución (Verde)
        return RGBColor(192, 0, 0)
    elif "↑" in parte_color:
        # Flecha hacia arriba = Aumento (rojo)
        return RGBColor(192, 0, 0)

    # Detectar por contexto del texto
    if ("sin cambios" in texto_completo or "sin cambio" in texto_completo) and (not "0" in texto_completo):
        # Sin cambio (Azul)
        return RGBColor(192, 0, 0)
    elif "disminuci" in texto_completo or "reducid" in texto_completo:
        # Disminución (Verde)
        return RGBColor(192, 0, 0)
    elif "aument" in texto_completo or "creci" in texto_completo:
        # Aumento (Rojo)
        return RGBColor(192, 0, 0)
    elif ("disminuci" in texto_completo or "reducid" in texto_completo) and ("0" in texto_completo):
        # Disminución (Verde)
        return RGBColor(0, 112, 192)
    elif "aument" in texto_completo or "creci" in texto_completo and ("0" in texto_completo):
        # Aumento (Rojo)
        return RGBColor(0, 112, 192)
    # ==== Nuevo ====
    elif ("sin cambios" in texto_completo or "sin cambio" in texto_completo) and "0" in texto_completo:
        # Sin cambio (Azul)
        return RGBColor(0, 112, 192)
    # ==== Nuevo ====
    
    # Si no se puede determinar por contexto, usar símbolo
    if "=" in parte_color:
        # Símbolo "=" normalmente indica aumento
        return RGBColor(0, 112, 192)
    
    # Fallback: Azul por defecto
    return RGBColor(0, 112, 192)

def reemplazar_simple(text: str, valores: dict) -> str:
    """
    Reemplazo simple de banderas (sin formato de color).
    """
    valores_norm = {}
    for k, v in valores.items():
        k_norm = k.lower()
        k_norm = re.sub(r'[\s\-\.]+', '_', k_norm)
        valores_norm[k_norm] = v
    
    def repl(m):
        bandera = m.group(1)
        bandera_norm = bandera.lower()
        bandera_norm = re.sub(r'[\s\-\.]+', '_', bandera_norm)
        valor_reemplazo = str(valores_norm.get(bandera_norm, m.group(0)))
        print(f"                     ✓ <<{bandera}>> → {valor_reemplazo}")
        return valor_reemplazo
    return re.sub(r"<<(.*?)>>", repl, text)
