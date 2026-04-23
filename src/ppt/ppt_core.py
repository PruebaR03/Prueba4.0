import os
import re
import pandas as pd
from datetime import datetime
from pptx import Presentation
from .ppt_utils import leer_hojas_excel, comparar_excel, buscar_bandera
from .ppt_bandera_handler import reemplazar_con_formato, reemplazar_simple

def cargar_resumen_cumplimiento(ruta_resumen_cumplimiento: str = None) -> dict:
    """
    Lee el archivo resumen_cumplimiento.xlsx y retorna un diccionario mapeado por métrica normalizada.

    Estructura de salida:
    {
        "crowdstrike": {
            "metrica": "Crowdstrike",
            "cumplimiento": "99.8%",
            "status": "Failed",
            "texto": "99.8% / Failed"
        }
    }
    """
    if not ruta_resumen_cumplimiento or not os.path.exists(ruta_resumen_cumplimiento):
        return {}

    try:
        df = pd.read_excel(ruta_resumen_cumplimiento)

        columnas_normalizadas = {str(c).strip().lower(): c for c in df.columns}

        col_metrica = columnas_normalizadas.get("métrica") or columnas_normalizadas.get("metrica")
        col_cumplimiento = (
            columnas_normalizadas.get("cumplimiento")
            or columnas_normalizadas.get("resultado")
        )
        col_status = (
            columnas_normalizadas.get("status")
            or columnas_normalizadas.get("estatus")
            or columnas_normalizadas.get("estado")
        )

        if not col_metrica or not col_cumplimiento or not col_status:
            print("   ⚠️  El archivo resumen_cumplimiento no tiene las columnas requeridas: Métrica, Resultado, Estatus")
            return {}

        resultado = {}

        for _, row in df.iterrows():
            metrica_original = "" if pd.isna(row[col_metrica]) else str(row[col_metrica]).strip()
            if not metrica_original:
                continue

            cumplimiento_raw = row[col_cumplimiento]
            status_raw = row[col_status]

            # Cumplimiento
            if pd.isna(cumplimiento_raw) or str(cumplimiento_raw).strip() == "":
                cumplimiento_texto = "N/A"
            else:
                try:
                    valor_num = float(cumplimiento_raw)
                    cumplimiento_texto = f"{valor_num * 100:.1f}%"
                except Exception:
                    cumplimiento_texto = str(cumplimiento_raw).strip()
                    cumplimiento_texto = re.sub(r'[\x00-\x1F]+', '', cumplimiento_texto)

            # Status
            if pd.isna(status_raw) or str(status_raw).strip() == "":
                status_texto = "N/A"
            else:
                status_texto = str(status_raw).strip()
                status_texto = re.sub(r'[\x00-\x1F]+', '', status_texto)

            metrica_key = re.sub(r'[\s\-\.]+', '_', metrica_original.strip().lower())

            resultado[metrica_key] = {
                "metrica": metrica_original,
                "cumplimiento": cumplimiento_texto,
                "status": status_texto,
                "texto": f"{cumplimiento_texto} / {status_texto}"
            }

        print(f"   ✅ {len(resultado)} métricas cargadas desde resumen_cumplimiento")
        return resultado

    except Exception as e:
        print(f"   ⚠️  No se pudo leer resumen_cumplimiento: {e}")
        return {}

def generar_ppt_comparativo(
    ruta_excel_actual: str,
    ruta_excel_anterior: str,
    ruta_plantilla: str,
    ruta_salida: str,
    ruta_instrucciones: str = None,
    ruta_resumen_cumplimiento: str = None
):
    """
    Genera un pptx comparando dos Excel (actual vs anterior).
    Soporta operaciones con banderas: <<archivo(operacion)>>
    """
    print("\n" + "═" * 70)
    print("📊 GENERADOR DE PRESENTACIÓN COMPARATIVA")
    print("═" * 70)
    
    print(f"\n📂 Leyendo archivos...")
    print(f"   Excel actual     : {ruta_excel_actual}")
    print(f"   Excel anterior   : {ruta_excel_anterior}")
    
    # Leer ambos Excel
    resumenes_actual = leer_hojas_excel(ruta_excel_actual)
    total_hojas_actual = len(resumenes_actual)
    print(f"   ✅ {total_hojas_actual} hoja(s) en Excel actual")
    
    resumenes_anterior = leer_hojas_excel(ruta_excel_anterior)
    total_hojas_anterior = len(resumenes_anterior)
    print(f"   ✅ {total_hojas_anterior} hoja(s) en Excel anterior")

    # Leer instrucciones para operaciones y categorías -- Nuevo de agrega ruta_resumen_cumplimiento
    categorias_por_hoja = {}
    instrucciones = None
    resumen_cumplimiento_map = cargar_resumen_cumplimiento(ruta_resumen_cumplimiento)
    print("CLAVES DE RESUMEN_CUMPLIMIENTO:")
    for k in resumen_cumplimiento_map.keys():
        print(f" - {k}")
    MAPEO_RESUMEN_CUMPLIMIENTO = {
        "dlp": "dlp_trellix",
        "netskope": "netskope_workstations",
        "qualys": "qualys_cmdb",
        "harmony": "harmony_instalado",
        "mtd": "mtd_sincronizado",
        "mtd_client_version": "mtd_version",
        "moviles": "moviles_compliant",
        "moviles_update": "moviles_update"
        
    }
    
    if ruta_instrucciones:
        from ..core import leer_instrucciones
        try:
            instrucciones = leer_instrucciones(ruta_instrucciones)
            for inst in instrucciones:
                hoja = inst.get("archivo")
                categoria = inst.get("categoria", "").strip()
                if hoja and categoria:
                    categorias_por_hoja[hoja.lower()] = categoria
            print(f"   ✅ {len(categorias_por_hoja)} categoría(s) cargada(s)")
            print(f"   ✅ Instrucciones cargadas para operaciones")
        except Exception as e:
            print(f"   ⚠️  No se pudo leer instrucciones: {e}")
    
    print(f"\n🔍 Comparando datos...")
    comparativo = comparar_excel(resumenes_actual, resumenes_anterior)
    fecha_actual = datetime.now().strftime("%d/%m/%Y")
    print(f"   📅 Fecha actual     : {fecha_actual}")
    print(f"   ✅ Comparación completada")

    print(f"\n🖼️  Abriendo plantilla PowerPoint...")
    print(f"   📁 {ruta_plantilla}")
    prs = Presentation(ruta_plantilla)
    print(f"   ✅ {len(prs.slides)} diapositiva(s) cargada(s)")

    print(f"\n" + "═" * 70)
    print("🎨 PROCESANDO DIAPOSITIVAS")
    print("═" * 70)
    
    diapositivas_procesadas = 0
    diapositivas_con_cambios = 0
    banderas_reemplazadas = 0

    for idx_slide, slide in enumerate(prs.slides, 1):
        banderas_encontradas = set()
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            banderas = buscar_bandera(shape.text)
            banderas_encontradas.update(banderas)
        
        if not banderas_encontradas:
            print(f"\n[{idx_slide}] '{slide.name}'")
            print(f"    ⊘ Estado: SIN BANDERAS (omitida)")
            continue

        print(f"\n[{idx_slide}] '{slide.name}'")
        print(f"    📌 Banderas: {', '.join(sorted(banderas_encontradas))}")

        slide_name = slide.name.lower() if slide.name else ""
        nombre_archivo = None
        categoria = None
        
        # Estrategia 1: Detectar por banderas
        if banderas_encontradas:
            print(f"    🔍 Detectando archivo...")
            
            for bandera_encontrada in banderas_encontradas:
                bandera_lower = bandera_encontrada.lower()
                # Extraer nombre si es operación: archivo(op) -> archivo
                op_match = re.match(r'^([^\(]+)\(', bandera_lower)
                if op_match:
                    bandera_lower = op_match.group(1).strip()
                
                for hoja_nombre in comparativo.keys():
                    hoja_norm = hoja_nombre.lower().replace(" ", "_")
                    
                    if hoja_norm == bandera_lower:
                        nombre_archivo = list(comparativo.keys())[0] if comparativo else None
                        print(f"    ✅ Archivo detectado por bandera")
                        break
                
                if nombre_archivo:
                    break
        
        # Estrategia 2: Por nombre de diapositiva
        if not nombre_archivo and slide_name:
            for archivo_key in comparativo.keys():
                archivo_norm = archivo_key.lower().replace(" ", "_").replace(".xlsx", "")
                if archivo_norm in slide_name:
                    nombre_archivo = archivo_key
                    print(f"    ✅ Archivo detectado por nombre")
                    break
        
        # Estrategia 3: Único archivo disponible
        if not nombre_archivo and comparativo:
            if len(comparativo) == 1:
                nombre_archivo = list(comparativo.keys())[0]
                print(f"    ✅ Único archivo disponible: {nombre_archivo}")
            else:
                print(f"    ⚠️  No se detectó archivo ({len(comparativo)} disponibles)")

        # Preparar valores
        valores = {}
        valores["fecha"] = fecha_actual
        nombre_archivo_actual = os.path.splitext(os.path.basename(ruta_excel_actual))[0]
        valores["titulo"] = nombre_archivo_actual
        
        # Detectar categoría
        categoria = None
        hojas = comparativo
        
        if banderas_encontradas:
            print(f"    🔍 Detectando categoría...")
            
            for bandera_encontrada in banderas_encontradas:
                bandera_lower = bandera_encontrada.lower()
                # Extraer nombre si es operación
                op_match = re.match(r'^([^\(]+)\(', bandera_lower)
                if op_match:
                    bandera_lower = op_match.group(1).strip()
                
                for hoja_nombre in hojas.keys():
                    hoja_norm = hoja_nombre.lower().replace(" ", "_")
                    
                    if hoja_norm == bandera_lower:
                        if hoja_nombre.lower() in categorias_por_hoja:
                            categoria = categorias_por_hoja[hoja_nombre.lower()]
                            print(f"    ✅ Categoría: {categoria}")
                            break
                
                if categoria:
                    break
        
        if not categoria and hojas:
            categoria = list(hojas.keys())[0]
            print(f"    ⓘ Categoría por defecto: {categoria}")
        
        valores["categoria"] = categoria if categoria else ""
        
        for key_resumen, info in resumen_cumplimiento_map.items():
            
            if f"{key_resumen}_cumplimiento" in valores:
                continue
            valores[f"{key_resumen}_cumplimiento"] = info["texto"]
            valores[f"{key_resumen}_status"] = info["status"]
            
            print(f"[SOLO RESUMEN] {key_resumen}_cumplimiento) -> {info['texto']}")

                # Alias de cumplimiento que no dependen directamente del nombre de la hoja -- Nuevo
        info_netskope_ws = resumen_cumplimiento_map.get("netskope_workstations")
        if info_netskope_ws:
            valores["netskope_workstations_cumplimiento"] = info_netskope_ws["texto"]
            valores["netskope_workstations_status"] = info_netskope_ws["status"]
        else:
            valores["netskope_workstations_cumplimiento"] = "N/A / N/A"
            valores["netskope_workstations_status"] = "N/A"

        info_netskope_pcv = resumen_cumplimiento_map.get("netskope_pc_virtuales")
        if info_netskope_pcv:
            valores["netskope_pc_virtuales_cumplimiento"] = info_netskope_pcv["texto"]
            valores["netskope_pc_virtuales_status"] = info_netskope_pcv["status"]
        else:
            valores["netskope_pc_virtuales_cumplimiento"] = "N/A / N/A"
            valores["netskope_pc_virtuales_status"] = "N/A"
        
        print(f"    ┌─────────────────────────────────")
        print(f"    │ 📋 VALORES A REEMPLAZAR:")
        print(f"    ├─────────────────────────────────")
        print(f"    │ Titulo    : {valores['titulo']}")
        print(f"    │ Categoria : {valores['categoria']}")
        print(f"    │ Fecha     : {valores['fecha']}")

        print("HOJAS EN COMPARATIVO:")
        for h in hojas.keys():
            print(f" - {h}")
        for hoja, datos in hojas.items():
            bandera = hoja.lower()
            bandera = re.sub(r'[\s\-\.]+', '_', bandera)

            print(f"hoja original: {hoja}")
            print(f"bandera generada: {bandera}")
            print(f"bandera cumplimiento generada: {bandera}_cumplimiento")
            
            valores[bandera] = datos['texto']
            valores[f'{bandera}_color'] = datos['color']

            # Nuevo: agregar bandera de cumplimiento/status por métrica
            metrica_key = re.sub(r'[\s\-\.]+', '_', hoja.strip().lower())
            metrica_key_resumen = MAPEO_RESUMEN_CUMPLIMIENTO.get(metrica_key, metrica_key)
            info_cumplimiento = resumen_cumplimiento_map.get(metrica_key_resumen)


            print(f"metrica_key original: {metrica_key}")
            print(f"metrica_key resumen : {metrica_key_resumen}")
            print(f"existe en resumen   : {'SI' if metrica_key_resumen in resumen_cumplimiento_map else 'NO'}")

            if info_cumplimiento:
                valores[f'{bandera}_cumplimiento'] = info_cumplimiento["texto"]
                valores[f'{bandera}_status'] = info_cumplimiento["status"]
            else:
                valores[f'{bandera}_cumplimiento'] = "N/A / N/A"
                valores[f'{bandera}_status'] = "N/A"

            color_rgb = datos['color']
            if hasattr(color_rgb, 'rgb'):
                color_hex = f"#{color_rgb.rgb:06X}"
            else:
                color_int = (color_rgb[0] << 16) | (color_rgb[1] << 8) | color_rgb[2]
                color_hex = f"#{color_int:06X}"

            print(f"    │ {bandera:<20} : {datos['texto']:<20} | {color_hex}")
            print(f"    │ {bandera + '_cumplimiento':<20} : {valores[f'{bandera}_cumplimiento']}")
        
        print(f"    ├─────────────────────────────────")
        print(f"    │ 🔄 REEMPLAZANDO BANDERAS:")
        print(f"    └─────────────────────────────────")
        
        for idx_shape, shape in enumerate(slide.shapes, 1):
            if not hasattr(shape, "text"):
                continue
            
            texto_original = shape.text
            banderas_en_shape = buscar_bandera(texto_original)
            if not banderas_en_shape:
                continue
            
            print(f"       📄 Cuadro {idx_shape}: {len(banderas_en_shape)} bandera(s)")
            
            if hasattr(shape, "text_frame"):
                try:
                    reemplazar_con_formato(
                        shape.text_frame, 
                        texto_original, 
                        valores, 
                        excel_actual=ruta_excel_actual,
                        excel_anterior=ruta_excel_anterior, 
                        instrucciones=instrucciones,
                        resumen_cumplimiento_map=resumen_cumplimiento_map
                    )
                    banderas_reemplazadas += len(banderas_en_shape)
                    print(f"          ✅ {len(banderas_en_shape)} reemplazada(s)")
                except Exception as e:
                    print(f"          ⚠️  Error: {str(e)[:40]}...")
                    nuevo_texto = reemplazar_simple(texto_original, valores)
                    if nuevo_texto != texto_original:
                        shape.text = nuevo_texto
                        banderas_reemplazadas += len(banderas_en_shape)
                        print(f"          ✅ {len(banderas_en_shape)} reemplazada(s) (fallback)")
            else:
                nuevo_texto = reemplazar_simple(texto_original, valores)
                if nuevo_texto != texto_original:
                    shape.text = nuevo_texto
                    banderas_reemplazadas += len(banderas_en_shape)
                    print(f"          ✅ {len(banderas_en_shape)} reemplazada(s) (simple)")
        
        diapositivas_con_cambios += 1
        diapositivas_procesadas += 1

    print(f"\n" + "═" * 70)
    print(f"📊 RESUMEN DE CAMBIOS")
    print("═" * 70)
    print(f"   📈 Diapositivas procesadas  : {diapositivas_procesadas}")
    print(f"   ✏️  Diapositivas con cambios : {diapositivas_con_cambios}")
    print(f"   🔀 Banderas reemplazadas    : {banderas_reemplazadas}")
    print("═" * 70)

    print(f"\n💾 Guardando archivo PowerPoint...")
    print(f"   📁 {ruta_salida}")
    prs.save(ruta_salida)
    print(f"   ✅ Archivo guardado exitosamente")
    
    print("\n" + "═" * 70)
    print("✅ PRESENTACION GENERADA COMPLETAMENTE")
    print("═" * 70 + "\n")
